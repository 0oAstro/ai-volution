
import streamlit as st
import openai
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import io
import os
from pathlib import Path
import re

from dotenv import load_dotenv

load_dotenv()

# Initialize OpenAI client
openai.api_key = os.getenv("OPENAI_API_KEY")
openai.base_url = os.getenv("OPENAI_API_URL")

def sanitize_category(name):
    """Clean category names for filesystem safety"""
    return re.sub(r'[\\/*?:"<>|]', "", name).replace(" ", "_").strip()[:50]

def parse_category(gpt_response):
    """Extract category name from GPT response with validation"""
    try:
        category_line = [line for line in gpt_response.split('\n') if line.lower().startswith('category:')][0]
        category = category_line.split(': ')[1].strip()
        return sanitize_category(category) if category else "Uncategorized"
    except:
        return "Uncategorized"

def extract_text(file_bytes, file_type):
    """Extract text from different file types"""
    text = ""
    try:
        if file_type == "pdf":
            pdf_reader = PdfReader(io.BytesIO(file_bytes))
            text = " ".join([page.extract_text() for page in pdf_reader.pages if page.extract_text()])
        elif file_type == "docx":
            doc = Document(io.BytesIO(file_bytes))
            text = "\n".join([para.text for para in doc.paragraphs if para.text])
        elif file_type == "pptx":
            prs = Presentation(io.BytesIO(file_bytes))
            text = " ".join([
                shape.text for slide in prs.slides 
                for shape in slide.shapes 
                if hasattr(shape, "text") and shape.text
            ])
        return text.strip()[:15000]  # Return first 15k characters
    except Exception as e:
        st.error(f"Error extracting text: {str(e)}")
        return None

def analyze_with_gpt(text, existing_categories=None):
    """Use OpenAI GPT-4 for document categorization"""
    system_message = """
# File Categorization System

You are an AI file categorization system. Your task is to automatically analyze and categorize files based on their content, name, and metadata.

## Core Categories
1.⁠ ⁠Finance (budgets, invoices, reports)
2.⁠ ⁠Legal (contracts, agreements, policies)
3.⁠ ⁠Project Management (plans, timelines, sprints)
4.⁠ ⁠Technical (code, documentation, APIs)
5.⁠ ⁠Marketing (campaigns, media, brands)
6.⁠ ⁠HR (employee docs, recruitment)
7.⁠ ⁠Presentations (slides, decks)
8.⁠ ⁠Personal (taxes, medical, family)

## Analysis Method
1.⁠ ⁠Check file extension/type
2.⁠ ⁠Scan filename for keywords
3.⁠ ⁠Analyze content patterns
4.⁠ ⁠Review metadata
5.⁠ ⁠Consider file location/context

## Response Format
For each file, respond ONLY in this format:

Category: [Specific category name]
Confidence: [X%]
Reason: [Brief technical explanation of categorization logic]

## Example Analysis
Input: "Q4_Revenue_Report_2024.xlsx"

Category: Finance
Confidence: 95%
Reason: Financial keywords detected, quarterly report pattern, spreadsheet format with revenue calculations

## Classification Rules
•⁠  ⁠Use highest confidence category
•⁠  ⁠Flag sensitive documents (HR, Legal)
•⁠  ⁠Consider file location context
•⁠  ⁠Pattern-match against known templates
•⁠  ⁠Check for department-specific markers

## Confidence Levels
•⁠  ⁠90%+ : Strong pattern match
•⁠  ⁠80-89%: Good indicators
•⁠  ⁠70-79%: Probable match
•⁠  ⁠<70%: Needs review

Never include additional commentary or explanations outside this format.
"""
    
    if existing_categories:
        system_message += f"\n\nExisting categories to consider: {', '.join(existing_categories)}"
    
    try:
        response = openai.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {"role": "system", "content": system_message},
                {"role": "user", "content": f"Document content:\n{text}"}
            ]
        )
        return response.choices[0].message.content
    except Exception as e:
        st.error(f"OpenAI API Error: {str(e)}")
        return None

def get_existing_categories():
    """Get list of existing categories from organized_docs directory"""
    base_dir = Path("organized_docs")
    if base_dir.exists():
        return [d.name for d in base_dir.iterdir() if d.is_dir()]
    return []

def process_files(uploaded_files):
    """Process all files and organize into categories"""
    results = []
    organized_files = {}
    existing_categories = get_existing_categories()
    
    for file in uploaded_files:
        file_bytes = file.getvalue()
        file_type = file.name.split(".")[-1].lower()
        
        with st.spinner(f"Processing {file.name}..."):
            text = extract_text(file_bytes, file_type)
            if not text:
                results.append((file.name, "Failed: Text extraction error", ""))
                continue
            
            gpt_response = analyze_with_gpt(text, existing_categories)
            if not gpt_response:
                results.append((file.name, "Failed: API error", ""))
                continue
            
            category = parse_category(gpt_response)
            sanitized_category = sanitize_category(category)
            
            # Save file to organized structure
            output_dir = Path("organized_docs") / sanitized_category
            output_dir.mkdir(parents=True, exist_ok=True)
            output_path = output_dir / file.name
            
            try:
                with open(output_path, "wb") as f:
                    f.write(file_bytes)
                organized_files.setdefault(sanitized_category, []).append(file.name)
                results.append((file.name, category, output_path))
                
                # Update existing categories if new one was created
                if sanitized_category not in existing_categories:
                    existing_categories.append(sanitized_category)
            except Exception as e:
                results.append((file.name, f"Failed: {str(e)}", ""))

    return results, organized_files

# Streamlit UI
st.title("📁 Smart Document Organizer")
st.write("Upload documents (PDF, DOCX, PPTX) for bulk processing and organization")

uploaded_files = st.file_uploader("Choose files", 
                                type=["pdf", "docx", "pptx"],
                                accept_multiple_files=True)

if uploaded_files:
    if st.button("🚀 Process All Files", help="Analyze and organize all uploaded files"):
        results, organized_files = process_files(uploaded_files)
        
        st.subheader("Processing Results")
        for file_name, category, path in results:
            with st.expander(f"{file_name} - {category.split(':')[0]}"):
                if "Failed" in category:
                    st.error(category)
                else:
                    st.success(f"Category: {category}")
                    st.code(f"Saved to: {path}")
        
        st.subheader("Organization Structure")
        for category, files in organized_files.items():
            st.write(f"📂 **{category}**")
            for f in files:
                st.write(f"└─ {f}")
            st.write("---")

else:
    st.info("Please upload documents to begin")